from flask import Flask, render_template, request, redirect, url_for
from pptx import Presentation
import difflib
import re

app = Flask(__name__)

ppt_text = ""


# ---------- EXTRACT PPT ----------
def extract_ppt_text(file):
    prs = Presentation(file)
    text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + " "

    return text.lower()


# ---------- LOGIN ----------
@app.route('/', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        return redirect(url_for('dashboard'))
    return render_template('index.html')


# ---------- DASHBOARD ----------
@app.route('/dashboard')
def dashboard():
    return render_template('dashboard.html')


# ---------- UPLOAD ----------
@app.route('/upload', methods=['POST'])
def upload():
    file = request.files.get('ppt')

    if file:
        file.save("uploaded.pptx")   # ✅ SAVE FILE permanently

    return redirect(url_for('dashboard'))


# ---------- ANALYZE ----------
@app.route('/analyze', methods=['POST'])
def analyze():
    import os

    # LOAD PPT
    if os.path.exists("uploaded.pptx"):
        ppt_text = extract_ppt_text("uploaded.pptx")
    else:
        ppt_text = ""

    spoken = request.form.get('text', '').lower()

    if spoken.strip() == "":
        return render_template('dashboard.html', error="No speech detected")

    # WORD LIST (for fillers, pauses, speed)
    spoken_list = re.findall(r'\b\w+\b', spoken.lower())

    # WORD SET (for accuracy)
    spoken_set = set(spoken_list)
    ppt_set = set(re.findall(r'\b\w+\b', ppt_text.lower()))

    # ACCURACY
    if spoken_set and ppt_set:
        common = spoken_set.intersection(ppt_set)
        accuracy = (len(common) / len(spoken_set)) * 100
    else:
        accuracy = 0

    # FILLERS
    filler_list = ["um", "uh", "like", "basically", "actually", "so", "and", "but"]
    fillers = sum(1 for word in spoken_list if word in filler_list)

    # PAUSES
    pauses = len(spoken_list) // 7

    # SPEED
    wpm = len(spoken_list)

    # FEEDBACK
    feedback = []

    if accuracy < 40:
        feedback.append("Follow PPT more closely")
    else:
        feedback.append("Good alignment")

    if fillers > 3:
        feedback.append("Reduce filler words")

    if wpm < 70:
        feedback.append("Speak faster")
    elif wpm > 150:
        feedback.append("Slow down")

    if pauses > 3:
        feedback.append("Too many pauses")

    return render_template(
        'dashboard.html',
        spoken=spoken,
        accuracy=round(accuracy, 2),
        fillers=fillers,
        pauses=pauses,
        wpm=wpm,
        feedback=feedback
    )
    if _name_ == "_main_": 
        app.run(debug=True)
